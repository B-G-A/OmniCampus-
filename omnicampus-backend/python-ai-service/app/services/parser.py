"""
OmniCampus AI Service — Document Parser.

Extracts raw text from PDF, DOCX, PPTX, and TXT files.  A single
dispatcher function `parse_document` routes to the correct extractor
based on the file_type argument.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


# ── Helpers ──────────────────────────────────────────────────────────────


def _clean_text(text: str) -> str:
    """Normalise whitespace and strip non‑printable / non‑unicode characters."""
    # Remove non-printable characters (keep newlines and tabs initially)
    text = re.sub(r"[^\x20-\x7E\n\t\r]", " ", text)
    # Collapse runs of whitespace on each line
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    # Collapse more than two consecutive blank lines into two
    cleaned: list[str] = []
    blank_count = 0
    for line in lines:
        if line == "":
            blank_count += 1
            if blank_count <= 2:
                cleaned.append(line)
        else:
            blank_count = 0
            cleaned.append(line)
    return "\n".join(cleaned).strip()


# ── Individual parsers ───────────────────────────────────────────────────


def parse_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pdfplumber (page by page)."""
    pages: list[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                pages.append(page_text)
    raw = "\n\n".join(pages)
    return _clean_text(raw)


def parse_docx(file_path: str) -> str:
    """Extract text from a DOCX file — paragraphs *and* table cells."""
    doc = DocxDocument(file_path)
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    return _clean_text("\n\n".join(parts))


def parse_pptx(file_path: str) -> str:
    """Extract text from a PPTX file — every text‑bearing shape on every slide."""
    prs = Presentation(file_path)
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    return _clean_text("\n\n".join(parts))


def parse_txt(file_path: str) -> str:
    """Read a plain‑text file."""
    content = Path(file_path).read_text(encoding="utf-8", errors="replace")
    return _clean_text(content)


# ── Dispatcher ───────────────────────────────────────────────────────────

_PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "txt": parse_txt,
}


def parse_document(file_path: str, file_type: str) -> str:
    """Route to the correct parser based on *file_type* (e.g. 'pdf').

    Raises ``ValueError`` for unsupported file types.
    """
    file_type = file_type.lower().strip().lstrip(".")
    parser = _PARSERS.get(file_type)
    if parser is None:
        raise ValueError(
            f"Unsupported file type '{file_type}'. "
            f"Supported types: {', '.join(_PARSERS)}"
        )
    logger.info("Parsing %s file: %s", file_type.upper(), file_path)
    text = parser(file_path)
    logger.info(
        "Extracted %d characters from %s",
        len(text),
        Path(file_path).name,
    )
    return text
